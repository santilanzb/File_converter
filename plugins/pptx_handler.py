"""
Concrete implementation of FileHandler for PowerPoint presentations (.pptx).

This module uses python-pptx to read and write PowerPoint presentations.
"""

from pptx import Presentation
from typing import List, Dict, Any

# Try to import from project structure, fall back to dummy classes
try:
    from ..core.exceptions import FileProcessingError
    from .base_handler import FileHandler
except ImportError:
    try:
        # Try absolute imports for installed package
        from core.exceptions import FileProcessingError
        from plugins.base_handler import FileHandler
    except ImportError:
        # For standalone testing, we'll define dummy classes.
        class FileProcessingError(Exception):
            def __init__(self, path, msg):
                super().__init__(f"{msg}: {path}")

        class FileHandler:
            def read(self, file_path: str): pass
            def write(self, file_path: str, data): pass




class PptxHandler(FileHandler):
    """
    Handles reading and writing of PowerPoint presentations (.pptx files).
    
    Extracts text from slides and creates presentations.
    """

    def read(self, file_path: str) -> List[Dict[str, Any]]:
        """
        Reads a PowerPoint presentation and extracts text content.

        Args:
            file_path (str): The path to the input PowerPoint presentation.

        Returns:
            List[Dict[str, Any]]: A list where each dictionary represents content
                                 with slide number and text information.

        Raises:
            FileProcessingError: If the file cannot be found, read, or processed.
        """
        try:
            presentation = Presentation(file_path)
            content_data = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_content = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_content += shape.text.strip() + "\n"
                content_data.append({
                    'type': 'slide',
                    'number': slide_num,
                    'content': slide_content.strip()
                })
            
            return content_data
            
        except FileNotFoundError:
            raise FileProcessingError(file_path, "PowerPoint presentation not found")
        except Exception as e:
            raise FileProcessingError(file_path, f"An error occurred while reading the PowerPoint presentation: {e}")

    def write(self, file_path: str, data: List[Dict[str, Any]]) -> None:
        """
        Writes data to a PowerPoint presentation.

        Args:
            file_path (str): The path to the output PowerPoint presentation.
            data (List[Dict[str, Any]]): The data to be written.

        Raises:
            FileProcessingError: If the file cannot be written.
            ValueError: If the data is invalid.
        """
        if not data:
            raise ValueError("Input data for PowerPoint presentation writing cannot be empty.")
        
        try:
            presentation = Presentation()
            
            for item in data:
                if isinstance(item, dict) and 'content' in item:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[5]) # Blank slide
                    text_box = slide.shapes.add_textbox(left=0, top=0, width=720, height=540)
                    text_frame = text_box.text_frame
                    text_frame.text = item['content']
                elif 'text' in item:
                    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                    text_box = slide.shapes.add_textbox(left=0, top=0, width=720, height=540)
                    text_frame = text_box.text_frame
                    text_frame.text = item['text']
                else:
                    # Handle non-dict items
                    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                    text_box = slide.shapes.add_textbox(left=0, top=0, width=720, height=540)
                    text_frame = text_box.text_frame
                    text_frame.text = str(item)
            
            presentation.save(file_path)
            
        except Exception as e:
            raise FileProcessingError(file_path, f"An error occurred while writing the PowerPoint presentation: {e}")
